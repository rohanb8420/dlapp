"""Utility for ingesting document metadata from an Excel training file into SQLite.

The script reads a spreadsheet (default: trainingdata.xlsx) containing
`filepath` and `businesscapability` columns, extracts light-weight metadata and
optionally textual content for each listed file, and persists the information in
`artifacts/dlm_reader.db`. A Streamlit UI is provided to explore the stored
records.
"""

import argparse
import csv
import hashlib
import io
import logging
import mimetypes
import re
import shutil
import sqlite3
import subprocess

import zipfile
import xml.etree.ElementTree as ET
from datetime import datetime
from pathlib import Path
from typing import Any, Dict, Iterable, List, Optional, Sequence

import pandas as pd

try:  # Optional UI dependency
    import streamlit as st
except Exception:  # pragma: no cover - environments without Streamlit support
    st = None

# Optional libraries for richer content extraction. The script falls back to
# empty strings when these packages are unavailable.
try:
    from docx import Document as DocxDocument
except Exception:  # pragma: no cover
    DocxDocument = None

try:
    from pptx import Presentation as PptxPresentation
except Exception:  # pragma: no cover
    PptxPresentation = None

try:
    import openpyxl
except Exception:  # pragma: no cover
    openpyxl = None

try:
    from PyPDF2 import PdfReader
except Exception:  # pragma: no cover
    PdfReader = None

try:  # Apache Tika (optional, requires Java) for hard-to-parse formats
    from tika import parser as tika_parser
except Exception:  # pragma: no cover
    tika_parser = None

LOGGER = logging.getLogger(__name__)
MAX_TEXT = 100_000
DB_PATH_DEFAULT = Path("artifacts") / "dlm_reader.db"

FILEPATH_COLUMN = "filepath"
FILEPATH_ALIASES: Sequence[str] = ("file_location", "file_path", "path")
CATEGORY_COLUMN = "businesscapability"
CATEGORY_ALIASES: Sequence[str] = (
    "business_category",
    "businessCategory",
    "BusinessCapability",
)

SCHEMA = """
CREATE TABLE IF NOT EXISTS files(
  file_id INTEGER PRIMARY KEY,
  path TEXT UNIQUE,
  folder TEXT,
  file_name TEXT,
  extension TEXT,
  mime_type TEXT,
  size_bytes INTEGER,
  created_ts TEXT,
  modified_ts TEXT,
  sha1 TEXT,
  exists_flag INTEGER,
  read_error TEXT
);
CREATE TABLE IF NOT EXISTS labels(
  file_id INTEGER,
  business_category TEXT,
  UNIQUE(file_id, business_category),
  FOREIGN KEY(file_id) REFERENCES files(file_id)
);
CREATE TABLE IF NOT EXISTS content(
  file_id INTEGER PRIMARY KEY,
  content_text TEXT,
  FOREIGN KEY(file_id) REFERENCES files(file_id)
);
CREATE INDEX IF NOT EXISTS idx_files_path ON files(path);
CREATE INDEX IF NOT EXISTS idx_labels_file ON labels(file_id);
"""

ASCII_RE = re.compile(rb"[\x20-\x7E]{4,}")


def ensure_db(db_path: Path) -> None:
    db_path.parent.mkdir(parents=True, exist_ok=True)
    con = sqlite3.connect(str(db_path))
    with con:
        con.executescript(SCHEMA)
    con.close()


def sha1_of_file(path: Path) -> Optional[str]:
    try:
        h = hashlib.sha1()
        with path.open("rb") as handle:
            for chunk in iter(lambda: handle.read(8192), b""):
                h.update(chunk)
        return h.hexdigest()
    except Exception:
        return None


def stat_path(path_str: str) -> Dict[str, Any]:
    path = Path(path_str)
    info = {
        "path": str(path),
        "folder": str(path.parent),
        "file_name": path.name,
        "extension": path.suffix.lower().lstrip("."),
        "mime_type": mimetypes.guess_type(str(path))[0] or "",
        "size_bytes": None,
        "created_ts": None,
        "modified_ts": None,
        "sha1": None,
        "exists_flag": 0,
        "read_error": None,
    }
    try:
        stats = path.stat()
        info.update(
            {
                "size_bytes": int(stats.st_size),
                "created_ts": datetime.fromtimestamp(stats.st_ctime).isoformat(),
                "modified_ts": datetime.fromtimestamp(stats.st_mtime).isoformat(),
                "sha1": sha1_of_file(path),
                "exists_flag": 1,
            }
        )
    except Exception as exc:
        info["read_error"] = f"stat_error: {exc}"
    return info


def _read_txt(path: Path) -> str:
    try:
        with path.open("r", encoding="utf-8", errors="ignore") as handle:
            text = handle.read(MAX_TEXT)
            return text[:MAX_TEXT]
    except Exception:
        return ""


def _read_csv(path: Path) -> str:
    try:
        rows: List[str] = []
        with path.open("r", encoding="utf-8", errors="ignore", newline="") as handle:
            reader = csv.reader(handle)
            for row in reader:
                rows.append(", ".join(cell.strip() for cell in row if cell))
                if len("\n".join(rows)) > MAX_TEXT:
                    break
        return "\n".join(rows)[:MAX_TEXT]
    except Exception:
        return ""


def _read_pdf(path: Path) -> str:
    if PdfReader is None:
        return ""
    try:
        reader = PdfReader(str(path))
        parts: List[str] = []
        for page in reader.pages:
            try:
                text = page.extract_text() or ""
            except Exception:
                text = ""
            if text:
                parts.append(text)
            if len("\n".join(parts)) > MAX_TEXT:
                break
        return "\n".join(parts)[:MAX_TEXT]
    except Exception:
        return ""


def _read_docx(path: Path) -> str:
    if DocxDocument is None:
        return ""
    try:
        document = DocxDocument(str(path))
        parts = [para.text for para in document.paragraphs if para.text]
        return "\n".join(parts)[:MAX_TEXT]
    except Exception:
        return ""


def _read_pptx(path: Path) -> str:
    if PptxPresentation is None:
        return ""
    try:
        presentation = PptxPresentation(str(path))
        pieces: List[str] = []
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    pieces.append(shape.text)
            if len("\n".join(pieces)) > MAX_TEXT:
                break
        return "\n".join(pieces)[:MAX_TEXT]
    except Exception:
        return ""


def _read_xlsx(path: Path) -> str:
    if openpyxl is None:
        return ""
    try:
        wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
        text_rows: List[str] = []
        for sheet in wb.worksheets:
            for row in sheet.iter_rows(values_only=True):
                row_vals = [str(cell) for cell in row if cell is not None]
                if row_vals:
                    text_rows.append(" ".join(row_vals))
                if len("\n".join(text_rows)) > MAX_TEXT:
                    break
            if len("\n".join(text_rows)) > MAX_TEXT:
                break
        return "\n".join(text_rows)[:MAX_TEXT]
    except Exception:
        return ""


def _text_from_dataframe(df: pd.DataFrame) -> str:
    if df.empty:
        return ""
    buffer = io.StringIO()
    df.to_csv(buffer, index=False)
    return buffer.getvalue()[:MAX_TEXT]


def _read_xls(path: Path) -> str:
    try:
        df = pd.read_excel(str(path), engine=None)
    except Exception:
        return ""
    return _text_from_dataframe(df)


def _read_twb(path: Path) -> str:
    try:
        tree = ET.parse(str(path))
        root = tree.getroot()
        texts: List[str] = []
        for elem in root.iter():
            if elem.text and elem.text.strip():
                texts.append(elem.text.strip())
            for attr in ("name", "caption", "label", "value"):
                value = elem.attrib.get(attr)
                if value:
                    texts.append(value)
            if len("\n".join(texts)) > MAX_TEXT:
                break
        return "\n".join(texts)[:MAX_TEXT]
    except Exception:
        return ""


def _read_twbx(path: Path) -> str:
    try:
        with zipfile.ZipFile(str(path)) as archive:
            members = [m for m in archive.namelist() if m.endswith(".twb")]
            if not members:
                return ""
            texts: List[str] = []
            with archive.open(members[0]) as handle:
                data = handle.read()
            root = ET.fromstring(data)
            for elem in root.iter():
                if elem.text and elem.text.strip():
                    texts.append(elem.text.strip())
                for attr in ("name", "caption", "label", "value"):
                    value = elem.attrib.get(attr)
                    if value:
                        texts.append(value)
                if len("\n".join(texts)) > MAX_TEXT:
                    break
            return "\n".join(texts)[:MAX_TEXT]
    except Exception:
        return ""


def _read_sas_dataset(path: Path) -> str:
    try:
        df = pd.read_sas(str(path), encoding="utf-8", format="sas7bdat")
    except ValueError:
        # Fall back to pandas auto-detection for alternate SAS binary formats.
        try:
            df = pd.read_sas(str(path))
        except Exception:
            return ""
    except Exception:
        return ""

    if isinstance(df, pd.DataFrame):
        return _text_from_dataframe(df)
    return ""


def _read_dll(path: Path) -> str:
    try:
        with path.open("rb") as handle:
            # Scan the first ~2MB for printable ASCII sequences to avoid huge payloads.
            data = handle.read(min(2_000_000, MAX_TEXT * 4))
    except Exception:
        return ""

    matches = ASCII_RE.findall(data)
    if not matches:
        return ""
    text = "\n".join(chunk.decode("ascii", errors="ignore") for chunk in matches)
    return text[:MAX_TEXT]


def _read_with_tika(path: Path) -> str:
    if tika_parser is None:
        return ""
    try:
        parsed = tika_parser.from_file(str(path))
        return (parsed.get("content") or "")[:MAX_TEXT]
    except Exception:
        return ""


def extract_text(path: Path, extension: str) -> str:
    ext = extension.lower()
    if ext in {"txt", "log", "md"}:
        return _read_txt(path)
    if ext in {"csv"}:
        return _read_csv(path)
    if ext == "pdf":
        text = _read_pdf(path)
        return text or _read_with_tika(path)
    if ext == "docx":
        text = _read_docx(path)
        return text or _read_with_tika(path)
    if ext == "pptx":
        text = _read_pptx(path)
        return text or _read_with_tika(path)
    if ext in {"xlsx", "xlsm"}:
        text = _read_xlsx(path)
        return text or _read_with_tika(path)
    if ext == "xls":
        text = _read_xls(path)
        return text or _read_with_tika(path)
    if ext == "twb":
        return _read_twb(path)
    if ext == "twbx":
        return _read_twbx(path)
    if ext == "sas":
        return _read_txt(path)
    if ext in {"sas7bdat", "sasbdat"}:
        text = _read_sas_dataset(path)
        return text or _read_with_tika(path)
    if ext == "dll":
        return _read_dll(path)
    if ext in {"one", "ppt", "doc", "msg"}:
        return _read_with_tika(path)
    return _read_with_tika(path)


def upsert_file(con: sqlite3.Connection, meta: Dict[str, Any]) -> int:
    with con:
        con.execute(
            """
            INSERT INTO files(path, folder, file_name, extension, mime_type, size_bytes,
                              created_ts, modified_ts, sha1, exists_flag, read_error)
            VALUES (:path, :folder, :file_name, :extension, :mime_type, :size_bytes,
                    :created_ts, :modified_ts, :sha1, :exists_flag, :read_error)
            ON CONFLICT(path) DO UPDATE SET
                folder=excluded.folder,
                file_name=excluded.file_name,
                extension=excluded.extension,
                mime_type=excluded.mime_type,
                size_bytes=excluded.size_bytes,
                created_ts=excluded.created_ts,
                modified_ts=excluded.modified_ts,
                sha1=excluded.sha1,
                exists_flag=excluded.exists_flag,
                read_error=excluded.read_error
            """,
            meta,
        )
        cursor = con.execute("SELECT file_id FROM files WHERE path = ?", (meta["path"],))
        row = cursor.fetchone()
        if not row:
            raise RuntimeError(f"Failed to retrieve file_id for {meta['path']}")
        return int(row[0])


def replace_label(con: sqlite3.Connection, file_id: int, label: str) -> None:
    with con:
        con.execute("DELETE FROM labels WHERE file_id = ?", (file_id,))
        if label:
            con.execute(
                "INSERT OR IGNORE INTO labels(file_id, business_category) VALUES (?, ?)",
                (file_id, label),
            )


def replace_content(con: sqlite3.Connection, file_id: int, text: str) -> None:
    with con:
        con.execute("DELETE FROM content WHERE file_id = ?", (file_id,))
        con.execute(
            "INSERT INTO content(file_id, content_text) VALUES (?, ?)",
            (file_id, text),
        )


def _resolve_column(df: pd.DataFrame, canonical: str, aliases: Sequence[str]) -> str:
    if canonical in df.columns:
        return canonical
    for name in aliases:
        if name in df.columns:
            return name
    raise KeyError(
        f"Could not find required column '{canonical}'. Tried aliases: {aliases}. "
        f"Available columns: {list(df.columns)}"
    )


def ingest_from_excel(
    excel_path: Path,
    db_path: Path,
    limit: Optional[int] = None,
) -> Dict[str, int]:
    ensure_db(db_path)
    df = pd.read_excel(excel_path)
    filepath_col = _resolve_column(df, FILEPATH_COLUMN, FILEPATH_ALIASES)
    category_col = _resolve_column(df, CATEGORY_COLUMN, CATEGORY_ALIASES)
    if limit is not None:
        df = df.head(limit)
    con = sqlite3.connect(str(db_path))
    stats = {"inserted": 0, "errors": 0, "empty_text": 0, "skipped_missing": 0}
    try:
        rows = df[[filepath_col, category_col]].itertuples(index=False, name=None)
        for filepath_value, category_value in rows:
            if filepath_value is None or str(filepath_value).strip() == "":
                stats["skipped_missing"] += 1
                continue
            path_str = str(filepath_value)
            meta = stat_path(path_str)
            try:
                file_id = upsert_file(con, meta)
                replace_label(
                    con,
                    file_id,
                    str(category_value) if category_value is not None else "",
                )
                text = ""
                if meta["exists_flag"]:
                    text = extract_text(Path(path_str), meta["extension"] or "")
                if not text:
                    stats["empty_text"] += 1
                replace_content(con, file_id, text)
                stats["inserted"] += 1
            except Exception as exc:
                LOGGER.exception("Failed to ingest %s", path_str)
                stats["errors"] += 1
    finally:
        con.close()
    return stats


def fetch_summary_rows(db_path: Path, limit: Optional[int] = None) -> List[Dict[str, Any]]:
    con = sqlite3.connect(str(db_path))
    con.row_factory = sqlite3.Row
    query = """
        SELECT f.file_id,
               f.file_name,
               f.folder,
               f.extension,
               f.mime_type,
               COALESCE(l.business_category, '') AS business_category,
               f.exists_flag,
               f.modified_ts,
               f.size_bytes
        FROM files f
        LEFT JOIN labels l ON l.file_id = f.file_id
        ORDER BY f.file_id
    """
    params: Iterable[Any] = ()
    if limit is not None:
        query += " LIMIT ?"
        params = (limit,)
    rows = [dict(row) for row in con.execute(query, params).fetchall()]
    con.close()
    return rows


def fetch_file_detail(db_path: Path, file_id: int) -> Optional[Dict[str, Any]]:
    con = sqlite3.connect(str(db_path))
    con.row_factory = sqlite3.Row
    try:
        file_row = con.execute("SELECT * FROM files WHERE file_id = ?", (file_id,)).fetchone()
        if not file_row:
            return None
        labels = [row[0] for row in con.execute(
            "SELECT business_category FROM labels WHERE file_id = ?", (file_id,)
        ).fetchall()]
        content_row = con.execute(
            "SELECT content_text FROM content WHERE file_id = ?", (file_id,)
        ).fetchone()
        return {
            "file": dict(file_row),
            "labels": labels,
            "content": content_row[0] if content_row else "",
        }
    finally:
        con.close()


def is_streamlit_runtime() -> bool:
    if st is None:
        return False
    try:
        from streamlit.runtime.scriptrunner import get_script_run_ctx
    except Exception:  # pragma: no cover
        return False
    return get_script_run_ctx() is not None


def run_streamlit_app(db_path: Path, default_excel: Path) -> None:
    if st is None:
        raise RuntimeError("Streamlit is not installed in this environment.")

    st.set_page_config(page_title="Document Metadata Browser", layout="wide")
    st.title("Document Metadata Browser")

    if "last_ingest_stats" not in st.session_state:
        st.session_state["last_ingest_stats"] = None
    if "last_ingest_excel" not in st.session_state:
        st.session_state["last_ingest_excel"] = str(default_excel)

    with st.sidebar:
        st.header("Ingestion")
        excel_input = st.text_input("Excel file", st.session_state["last_ingest_excel"])
        limit_value = st.number_input(
            "Row limit (0 = all)", min_value=0, step=1, value=0, format="%d"
        )
        run_ingest = st.button("Run ingestion")
        if run_ingest:
            excel_path = Path(excel_input)
            if not excel_path.exists():
                st.error(f"Excel file not found: {excel_path}")
            else:
                limit = limit_value if limit_value > 0 else None
                with st.spinner("Ingesting files..."):
                    stats = ingest_from_excel(excel_path, db_path, limit=limit)
                st.session_state["last_ingest_stats"] = stats
                st.session_state["last_ingest_excel"] = excel_input
                st.success(
                    f"Ingested {stats['inserted']} rows | "
                    f"skipped_missing={stats['skipped_missing']} empty_text={stats['empty_text']} errors={stats['errors']}"
                )

        st.markdown("---")
        st.subheader("Database")
        st.write(f"Database path: `{db_path}`")

    stats = st.session_state.get("last_ingest_stats")
    if stats:
        st.toast(
            f"Last ingest: inserted={stats['inserted']} | skipped_missing={stats['skipped_missing']} | "
            f"empty_text={stats['empty_text']} | errors={stats['errors']}",
            icon="?" if stats.get("errors", 0) == 0 else "??",
        )

    rows = fetch_summary_rows(db_path)
    if not rows:
        st.info("No files in the database yet. Run an ingestion to populate the table.")
        return

    df_summary = pd.DataFrame(rows)
    st.subheader("Files")
    st.dataframe(
        df_summary,
        use_container_width=True,
        hide_index=True,
    )

    file_ids = df_summary["file_id"].tolist()
    default_file = file_ids[0] if file_ids else None
    selected_id = st.selectbox(
        "Select file", file_ids, index=0 if default_file is not None else None
    )

    if selected_id is None:
        st.warning("Select a file to view its metadata.")
        return

    detail = fetch_file_detail(db_path, int(selected_id))
    if not detail:
        st.warning("No details found for the selected file.")
        return

    file_meta = detail["file"]
    st.subheader("Metadata")
    st.json(file_meta)

    labels = detail["labels"] or ["(unlabeled)"]
    st.markdown("**Business Capability:** " + ", ".join(labels))

    content_preview = detail["content"][:2_000]
    st.subheader("Content Preview")
    st.text_area(
        "Extracted text (first 2,000 characters)",
        value=content_preview or "(no text extracted)",
        height=300,
    )


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(
        description="Ingest document metadata and optionally launch a Streamlit browser UI."
    )
    parser.add_argument(
        "--excel",
        default="trainingdata.xlsx",
        type=Path,
        help="Path to the Excel spreadsheet containing filepath/businesscapability columns.",
    )
    parser.add_argument(
        "--db",
        default=str(DB_PATH_DEFAULT),
        type=Path,
        help="Destination SQLite database path.",
    )
    parser.add_argument(
        "--limit",
        type=int,
        default=None,
        help="Optional limit on the number of rows to ingest from Excel.",
    )
    parser.add_argument(
        "--skip-ingest",
        action="store_true",
        help="Skip ingestion and only open the browser UI.",
    )
    parser.add_argument(
        "--no-ui",
        action="store_true",
        help="Run ingestion only and do not launch the UI.",
    )
    parser.add_argument(
        "--log-level",
        default="INFO",
        choices=["CRITICAL", "ERROR", "WARNING", "INFO", "DEBUG"],
        help="Logging verbosity for the script.",
    )
    args, unknown = parser.parse_known_args()
    if unknown:
        LOGGER.debug("Ignoring unknown arguments: %s", unknown)
    return args


def configure_logging(level: str) -> None:
    logging.basicConfig(
        level=getattr(logging, level.upper(), logging.INFO),
        format="%(asctime)s | %(levelname)s | %(name)s | %(message)s",
    )


def main() -> None:
    args = parse_args()
    configure_logging(args.log_level)

    db_path = Path(args.db)
    ensure_db(db_path)

    running_in_streamlit = is_streamlit_runtime()
    ingest_stats = None

    if not args.skip_ingest:
        excel_path = Path(args.excel)
        if running_in_streamlit and st is not None:
            if "initial_ingest_done" not in st.session_state:
                if excel_path.exists():
                    ingest_stats = ingest_from_excel(
                        excel_path, db_path, limit=args.limit
                    )
                    st.session_state["initial_ingest_stats"] = ingest_stats
                else:
                    LOGGER.warning("Excel file not found: %s", excel_path)
                st.session_state["initial_ingest_done"] = True
        else:
            if excel_path.exists():
                ingest_stats = ingest_from_excel(excel_path, db_path, limit=args.limit)
                LOGGER.info(
                    "Ingest complete | inserted=%s skipped_missing=%s empty_text=%s errors=%s",
                    ingest_stats["inserted"],
                    ingest_stats["skipped_missing"],
                    ingest_stats["empty_text"],
                    ingest_stats["errors"],
                )
            else:
                LOGGER.warning("Excel file not found: %s", excel_path)
    else:
        LOGGER.info("Skipping ingestion step as requested.")

    if args.no_ui:
        return

    if running_in_streamlit:
        run_streamlit_app(db_path, Path(args.excel))
        return

    if st is None:
        LOGGER.error(
            "Streamlit is not installed. Install it via `pip install streamlit` or re-run with --no-ui."
        )
        return

    streamlit_cli = shutil.which("streamlit")
    if not streamlit_cli:
        LOGGER.error(
            "Streamlit CLI not found in PATH. Install Streamlit or adjust your PATH environment."
        )
        return

    cmd = [
        streamlit_cli,
        "run",
        str(Path(__file__).resolve()),
        "--",
        "--excel",
        str(Path(args.excel).resolve()),
        "--db",
        str(db_path.resolve()),
    ]
    if args.limit is not None:
        cmd.extend(["--limit", str(args.limit)])
    if args.skip_ingest:
        cmd.append("--skip-ingest")
    if args.log_level:
        cmd.extend(["--log-level", args.log_level])

    LOGGER.info("Launching Streamlit app...")
    subprocess.run(cmd, check=False)


if __name__ == "__main__":
    main()


